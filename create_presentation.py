"""
Скрипт генерации PPTX презентации приложения AutoParser.
Создает современную стильную презентацию 16:9 с местами для скриншотов.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    
    # Формат 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Цветовая палитра
    COLOR_BG = RGBColor(0x1E, 0x22, 0x2A)        # Тёмно-серый задний фон
    COLOR_CARD = RGBColor(0x28, 0x2C, 0x34)      # Цвет карточек
    COLOR_BLUE = RGBColor(0x2F, 0x80, 0xED)      # Синий акцент
    COLOR_GREEN = RGBColor(0x27, 0xAE, 0x60)     # Зелёный акцент
    COLOR_TEXT = RGBColor(0xF0, 0xF2, 0xF5)      # Белый текст
    COLOR_MUTED = RGBColor(0xAA, 0xB2, 0xC0)     # Серый текст
    COLOR_PLACEHOLDER_BG = RGBColor(0x18, 0x1B, 0x20) # Фон рамки для скриншота

    blank_layout = prs.slide_layouts[6]  # Пустой слайд

    def set_slide_background(slide):
        """Установка тёмного фона слайда."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG

    def add_header(slide, title_text, subtitle_text):
        """Добавление стильного заголовка."""
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(1.0))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT
        
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(14)
        p2.font.color.rgb = COLOR_MUTED

    def add_screenshot_placeholder(slide, left, top, width, height, label="МЕСТО ДЛЯ СКРИНШОТА ИНТЕРФЕЙСА"):
        """Добавление аккуратного блока-плейсхолдера для скриншота."""
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_PLACEHOLDER_BG
        shape.line.color.rgb = COLOR_BLUE
        shape.line.width = Pt(2)
        
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = f"📷 {label}\n(Вставьте скриншот сюда)"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_MUTED

    # =========================================================================
    # СЛАЙД 1: Титульный слайд
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1)

    # Декоративный карточка по центру
    card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.2), Inches(10.333), Inches(5.1))
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_CARD
    card.line.color.rgb = COLOR_BLUE
    card.line.width = Pt(1.5)

    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.8)
    tf.margin_left = Inches(0.8)

    p1 = tf.paragraphs[0]
    p1.text = "🚗 AutoParser"
    p1.font.size = Pt(44)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TEXT

    p2 = tf.add_paragraph()
    p2.text = "Автономная система сбора и анализа данных автомобильных площадок и аукционов"
    p2.font.size = Pt(20)
    p2.font.color.rgb = COLOR_BLUE

    p3 = tf.add_paragraph()
    p3.text = "\n• Автоматический поиск и фильтрация объявлений (ОАЭ, Европа, США, Корея, Китай)\n• Встроенный интерактивный интерфейс (GUI) без необходимости сторонних СУБД\n• Быстрый экспорт результатов в Excel, CSV и JSON"
    p3.font.size = Pt(15)
    p3.font.color.rgb = COLOR_MUTED

    # =========================================================================
    # СЛАЙД 2: Обзор приложения и Главный экран GUI
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_header(slide2, "1. Главный экран приложения (GUI)", "Удобное десктопное приложение на базе CustomTkinter")

    # Текстовое описание слева
    info_box = slide2.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(4.5), Inches(5.2))
    tf = info_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "💡 Ключевые элементы:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_GREEN

    items = [
        ("Современный Dark Theme", "Эргономичный тёмный интерфейс, адаптированный под macOS и Windows."),
        ("Автономная работа", "Не требует настройки базы данных PostgreSQL — работает из коробки."),
        ("Управление в 1 клик", "Быстрый запуск поиска, выбор параметров и открытие готовых файлов."),
        ("Журнал событий (Logs)", "Отображение хода выполнения парсинга в реальном времени.")
    ]

    for title, desc in items:
        p_t = tf.add_paragraph()
        p_t.text = f"• {title}"
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_TEXT
        
        p_d = tf.add_paragraph()
        p_d.text = f"  {desc}"
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = COLOR_MUTED

    # Скриншот справа
    add_screenshot_placeholder(slide2, Inches(5.6), Inches(1.5), Inches(6.9), Inches(5.3), "Главное окно приложения")

    # =========================================================================
    # СЛАЙД 3: Поиск и фильтрация
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_header(slide3, "2. Настройка поиска и фильтрация", "Точный подбор автомобилей по марке, моделям и годам")

    # Скриншот слева
    add_screenshot_placeholder(slide3, Inches(0.8), Inches(1.5), Inches(6.9), Inches(5.3), "Панель фильтров поиска")

    # Описание справа
    info_box = slide3.shapes.add_textbox(Inches(8.0), Inches(1.5), Inches(4.5), Inches(5.2))
    tf = info_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "🎯 Параметры поиска:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLUE

    items = [
        ("Марка / Модель", "Поиск конкретных моделей (например: Mitsubishi Outlander, Toyota Land Cruiser)."),
        ("Диапазон лет", "Фильтрация по годам выпуска (например: 2022–2026 гг.)."),
        ("Лимит страниц", "Указание глубины сканирования (1, 3, 5+ страниц)."),
        ("Headless режим", "Запуск браузера в фоновом режиме без отображения окон.")
    ]

    for title, desc in items:
        p_t = tf.add_paragraph()
        p_t.text = f"• {title}"
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_TEXT
        
        p_d = tf.add_paragraph()
        p_d.text = f"  {desc}"
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = COLOR_MUTED

    # =========================================================================
    # СЛАЙД 4: Интерактивная таблица результатов
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4)
    add_header(slide4, "3. Интерактивные результаты", "Мгновенный просмотр собранных данных и экспорт")

    info_box = slide4.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(4.5), Inches(5.2))
    tf = info_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "📊 Возможности таблицы:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_GREEN

    items = [
        ("19 колонок данных", "Марка, модель, цена в AED, пробег, год, топливо, КПП, город."),
        ("Двойной клик", "Быстрый переход на страницу объявления в браузере по двойному клику."),
        ("Экспорт в Excel / CSV", "Сохранение структурированных отчётов в 1 клик."),
        ("Сводная статистика", "Отображение общего количества найденных вариантов.")
    ]

    for title, desc in items:
        p_t = tf.add_paragraph()
        p_t.text = f"• {title}"
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_TEXT
        
        p_d = tf.add_paragraph()
        p_d.text = f"  {desc}"
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = COLOR_MUTED

    add_screenshot_placeholder(slide4, Inches(5.6), Inches(1.5), Inches(6.9), Inches(5.3), "Таблица результатов поиска")

    # =========================================================================
    # СЛАЙД 5: Архитектура и поддерживаемые площадки
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5)
    add_header(slide5, "4. Архитектура и поддерживаемые площадки", "Модульная структура для масштабирования на другие рынки")

    # 4 карточки площадок
    platforms = [
        ("🇦🇪 DubiCars (ОАЭ)", "Парсинг цен в AED, спецификаций и дилеров. Активен и полностью протестирован.", COLOR_GREEN),
        ("🇪🇺 AutoScout24 (Европа)", "Поддержка фильтрации по странам (Германия, Франция, Италия, Нидерланды).", COLOR_BLUE),
        ("🇺🇸 BidCars (США)", "Парсинг лотов автоаукционов, конвертация миль в км, извлечение VIN.", COLOR_BLUE),
        ("🇨🇳 Che168 & Autohome (Китай)", "Поддержка дешифровки кастомных WebFont шрифтов (fontTools) для цен.", COLOR_BLUE),
    ]

    coords = [
        (Inches(0.8), Inches(1.6)),
        (Inches(6.8), Inches(1.6)),
        (Inches(0.8), Inches(4.3)),
        (Inches(6.8), Inches(4.3)),
    ]

    for idx, (p_title, p_desc, p_color) in enumerate(platforms):
        x, y = coords[idx]
        card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.7), Inches(2.4))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = p_color
        card.line.width = Pt(1.5)

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = Inches(0.2)
        
        p = tf.paragraphs[0]
        p.text = p_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT

        p2 = tf.add_paragraph()
        p2.text = f"\n{p_desc}"
        p2.font.size = Pt(13)
        p2.font.color.rgb = COLOR_MUTED

    # =========================================================================
    # СЛАЙД 6: Итоги и планы развития
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6)
    add_header(slide6, "5. Заключение и дальнейшие шаги", "Готово к практическому использованию")

    card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.5), Inches(10.333), Inches(5.1))
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_CARD
    card.line.color.rgb = COLOR_GREEN
    card.line.width = Pt(1.5)

    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.5)
    tf.margin_left = Inches(0.6)

    p1 = tf.paragraphs[0]
    p1.text = "✅ Текущий статус проекта:"
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_GREEN

    p2 = tf.add_paragraph()
    p2.text = "• Полностью рабочий парсер DubiCars с графическим интерфейсом.\n• Проверено извлечение данных по реальным запросам (Mitsubishi Outlander 2022-2026).\n• Весь код загружен в GitHub репозиторий.\n"
    p2.font.size = Pt(14)
    p2.font.color.rgb = COLOR_TEXT

    p3 = tf.add_paragraph()
    p3.text = "🚀 План развития:"
    p3.font.size = Pt(20)
    p3.font.bold = True
    p3.font.color.rgb = COLOR_BLUE

    p4 = tf.add_paragraph()
    p4.text = "• Подключение отображения остальных 5 площадок в едином GUI интерфейсе.\n• Добавление графиков аналитики цен по годам и пробегу.\n• Автоматическое уведомление в Telegram о появлении выгодных объявлений."
    p4.font.size = Pt(14)
    p4.font.color.rgb = COLOR_MUTED

    output_path = "/Users/zakhar/Desktop/AutoParser/AutoParser_Presentation.pptx"
    prs.save(output_path)
    print(f"Презентация создана: {output_path}")

if __name__ == "__main__":
    create_deck()
